"""
知识库文档导入：PDF / PPTX / TXT / 图片 → 多模态识别 → 按 00-18 模块体系结构化草稿。

管线（异步任务，独立轻量 worker 线程）：
  render    文件 → 页面负载（PDF 逐页渲染成图；PPTX 逐页取文本+内嵌图；图片直用；TXT 直读）
  recognize 图片页 → qwen3.7-plus（原生多模态，经中继）转录为 Markdown 文本
  extract   全文 → LLM 按模块体系结构化（每模块必须附原文引用 quotes）
  verify    确定性核验：quotes 与原文做归一化子串匹配，不再花 LLM 调用

产出是"草稿"（draft_json），必须经前端审核逐模块采纳后才合并进知识库——
知识库是准确率校验的标准答案，绝不让 LLM 直出裁判标准。
"""
from __future__ import annotations

import base64
import json
import queue
import re
import sqlite3
import threading
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from pathlib import Path
from typing import Any
from uuid import uuid4

import yaml

from ..core.paths import CONFIG_DIR, DATA_DIR

KB_IMPORTS_DB = DATA_DIR / "jobs.sqlite"
KB_IMPORTS_DIR = DATA_DIR / "kb_imports"

MAX_FILE_BYTES = 20 * 1024 * 1024
MAX_PAGES = 40
RECOGNIZE_CONCURRENCY = 5
VISION_MODEL_DEFAULT = "qwen3.7-plus"

ALLOWED_EXTS = {".pdf", ".pptx", ".txt", ".md", ".png", ".jpg", ".jpeg", ".webp"}

# 与 docs/products_db / 06_build_knowledge_base.py 一致的模块体系。
# 文档导入只面向"事实性"模块；表达库/版本记录类（12-14/17/18）不由文档抽取。
MODULE_CATALOG = {
    "01": "产品身份信息",
    "02": "官方权威资料",
    "03": "产品定位与核心价值",
    "04": "核心适用场景",
    "05": "次相关场景",
    "06": "不适用禁用谨慎场景",
    "07": "作用机制与原理",
    "08": "使用方法与操作规范",
    "09": "产品差异与产品矩阵",
    "10": "同类产品与竞品关系",
    "11": "关联联合使用",
    "15": "FAQ知识库",
    "16": "合规与风险提示",
}

EXTRACT_PROMPT = """你是药品知识库整理专家。请把下面的产品资料原文，按指定模块体系整理成知识库条目。

## 模块体系（只能使用这些编号）
{catalog}

## 铁律
1. 只收录原文中明确陈述的事实，禁止补充任何原文没有的"常识"或推测
2. 每个模块的 quotes 字段必须是原文的逐字摘录（每条 10-50 字），证明该模块内容有依据
3. 原文没有涉及的模块直接省略，不要编造
4. text 用简洁的要点式表述，保留剂量、频次、禁忌等关键数字和限定词

## 输出格式
严格返回 JSON 对象：
{{"modules": {{"04": {{"name": "核心适用场景", "text": "整理后的知识文本", "quotes": ["原文摘录1", "原文摘录2"]}}}}}}

## 产品
{product}

## 资料原文
{source}

只返回 JSON，不要任何其他文本。"""

RECOGNIZE_PROMPT = (
    "请把这页文档的全部文字内容完整转录为 Markdown（保留标题层级、表格、列表结构；"
    "忽略页眉页脚页码和纯装饰元素）。只输出转录内容，不要任何解释。"
)

_QUEUE: "queue.Queue[str]" = queue.Queue()
_WORKER_STARTED = False
_WORKER_LOCK = threading.Lock()


def _connect() -> sqlite3.Connection:
    KB_IMPORTS_DB.parent.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(KB_IMPORTS_DB)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    return conn


def init_db() -> None:
    with _connect() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS kb_imports (
                import_id TEXT PRIMARY KEY,
                username TEXT NOT NULL,
                product_key TEXT NOT NULL,
                scope TEXT NOT NULL DEFAULT 'user',
                filename TEXT NOT NULL,
                file_ext TEXT NOT NULL,
                status TEXT NOT NULL DEFAULT 'queued',
                stage TEXT,
                page_count INTEGER,
                pages_done INTEGER DEFAULT 0,
                error TEXT,
                draft_json TEXT,
                source_chars INTEGER,
                created_at TEXT NOT NULL,
                finished_at TEXT,
                applied_at TEXT
            )
            """
        )


def _now() -> str:
    return datetime.now().isoformat(timespec="seconds")


def _update(import_id: str, **updates: Any) -> None:
    keys = ", ".join(f"{k} = ?" for k in updates)
    with _connect() as conn:
        conn.execute(
            f"UPDATE kb_imports SET {keys} WHERE import_id = ?", (*updates.values(), import_id)
        )


def get_import(import_id: str) -> dict[str, Any] | None:
    with _connect() as conn:
        row = conn.execute(
            "SELECT * FROM kb_imports WHERE import_id = ?", (import_id,)
        ).fetchone()
    return dict(row) if row else None


def list_imports(username: str | None = None) -> list[dict[str, Any]]:
    with _connect() as conn:
        if username is None:
            rows = conn.execute(
                "SELECT * FROM kb_imports ORDER BY created_at DESC LIMIT 100"
            ).fetchall()
        else:
            rows = conn.execute(
                "SELECT * FROM kb_imports WHERE username = ? ORDER BY created_at DESC LIMIT 100",
                (username,),
            ).fetchall()
    items = [dict(row) for row in rows]
    for item in items:
        item.pop("draft_json", None)  # 列表接口不带大字段
    return items


def create_import(
    username: str,
    role: str,
    product_key: str,
    scope: str,
    filename: str,
    content: bytes,
) -> dict[str, Any]:
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXTS:
        raise ValueError(f"不支持的文件类型 {ext}（支持 PDF / PPTX / TXT / MD / 图片）")
    if len(content) > MAX_FILE_BYTES:
        raise ValueError("文件超过 20MB")
    if not product_key.strip():
        raise ValueError("请先选择/填写产品")
    if scope == "global" and role != "admin":
        scope = "user"
    import_id = f"kbi_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid4().hex[:6]}"
    workdir = KB_IMPORTS_DIR / import_id
    workdir.mkdir(parents=True, exist_ok=True)
    (workdir / f"source{ext}").write_bytes(content)
    with _connect() as conn:
        conn.execute(
            """
            INSERT INTO kb_imports (
                import_id, username, product_key, scope, filename, file_ext,
                status, created_at
            ) VALUES (?, ?, ?, ?, ?, ?, 'queued', ?)
            """,
            (import_id, username, product_key.strip(), scope, filename, ext, _now()),
        )
    ensure_worker()
    _QUEUE.put(import_id)
    return get_import(import_id)


def ensure_worker() -> None:
    global _WORKER_STARTED
    with _WORKER_LOCK:
        if _WORKER_STARTED:
            return
        _WORKER_STARTED = True
        with _connect() as conn:
            rows = conn.execute(
                "SELECT import_id FROM kb_imports WHERE status IN ('queued', 'running') ORDER BY created_at"
            ).fetchall()
        for row in rows:
            _QUEUE.put(row["import_id"])
        threading.Thread(target=_worker_loop, daemon=True).start()


def _worker_loop() -> None:
    while True:
        import_id = _QUEUE.get()
        try:
            _run_import(import_id)
        except Exception as exc:  # noqa: BLE001 —— worker 不允许死
            _update(import_id, status="failed", error=str(exc)[:500], finished_at=_now())
        finally:
            _QUEUE.task_done()


# ---------------------------------------------------------------------------
# LLM 客户端（OpenAI 兼容：优先中继，回退 dashscope 直连）
# ---------------------------------------------------------------------------


def _vision_client():
    from openai import OpenAI

    from utils.api_clients import resolve_relay

    keys = yaml.safe_load((CONFIG_DIR / "api_keys.yaml").read_text(encoding="utf-8")) or {}
    models_cfg = yaml.safe_load((CONFIG_DIR / "models.yaml").read_text(encoding="utf-8")) or {}
    relay = resolve_relay(models_cfg, keys)
    if relay:
        return OpenAI(base_url=relay["base_url"], api_key=relay["api_key"], timeout=180), "relay"
    qwen_key = (keys.get("qwen") or {}).get("api_key", "")
    if qwen_key and qwen_key != "sk-xxx":
        return (
            OpenAI(
                base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
                api_key=qwen_key,
                timeout=180,
            ),
            "direct",
        )
    raise RuntimeError("中继和千问直连都未配置，无法进行多模态识别")


def _vision_model() -> str:
    import os

    return os.environ.get("GEO_VISION_MODEL") or VISION_MODEL_DEFAULT


# ---------------------------------------------------------------------------
# render：文件 → 页面负载
# ---------------------------------------------------------------------------


def _render_pages(path: Path, ext: str) -> list[dict[str, Any]]:
    """返回页面负载列表：{"kind": "image", "b64": ..., "mime": ...} 或 {"kind": "text", "text": ...}"""
    if ext in (".txt", ".md"):
        raw = path.read_bytes()
        for encoding in ("utf-8-sig", "utf-8", "gbk"):
            try:
                return [{"kind": "text", "text": raw.decode(encoding)}]
            except UnicodeDecodeError:
                continue
        return [{"kind": "text", "text": raw.decode("utf-8", errors="replace")}]

    if ext in (".png", ".jpg", ".jpeg", ".webp"):
        mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg", "webp": "image/webp"}[ext[1:]]
        return [{"kind": "image", "b64": base64.b64encode(path.read_bytes()).decode(), "mime": mime}]

    if ext == ".pdf":
        import fitz

        pages = []
        with fitz.open(path) as doc:
            if doc.page_count > MAX_PAGES:
                raise ValueError(f"PDF 共 {doc.page_count} 页，超过 {MAX_PAGES} 页上限，请拆分后上传")
            for page in doc:
                pix = page.get_pixmap(dpi=140)
                pages.append({
                    "kind": "image",
                    "b64": base64.b64encode(pix.tobytes("png")).decode(),
                    "mime": "image/png",
                })
        return pages

    if ext == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(path))
        if len(prs.slides) > MAX_PAGES:
            raise ValueError(f"PPT 共 {len(prs.slides)} 页，超过 {MAX_PAGES} 页上限，请拆分后上传")
        pages = []
        for slide in prs.slides:
            texts = []
            images = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
                if shape.shape_type == 13 and getattr(shape, "image", None):  # PICTURE
                    blob = shape.image.blob
                    if blob and len(blob) < 5 * 1024 * 1024:
                        images.append({
                            "b64": base64.b64encode(blob).decode(),
                            "mime": shape.image.content_type or "image/png",
                        })
            page: dict[str, Any] = {"kind": "slide", "text": "\n".join(texts), "images": images[:3]}
            pages.append(page)
        return pages

    raise ValueError(f"未实现的类型: {ext}")


# ---------------------------------------------------------------------------
# recognize / extract / verify
# ---------------------------------------------------------------------------


def _recognize_page(client, model: str, page: dict[str, Any]) -> str:
    if page["kind"] == "text":
        return page["text"]
    content: list[dict[str, Any]] = []
    if page["kind"] == "image":
        content.append({
            "type": "image_url",
            "image_url": {"url": f"data:{page['mime']};base64,{page['b64']}"},
        })
        content.append({"type": "text", "text": RECOGNIZE_PROMPT})
    else:  # slide：文本 + 内嵌图片一起识别
        for image in page.get("images", []):
            content.append({
                "type": "image_url",
                "image_url": {"url": f"data:{image['mime']};base64,{image['b64']}"},
            })
        slide_text = page.get("text", "")
        if page.get("images"):
            content.append({
                "type": "text",
                "text": f"这是一页 PPT。文本框内容如下：\n{slide_text}\n\n"
                        f"请结合上面的图片，把这页 PPT 的全部信息转录为 Markdown。只输出转录内容。",
            })
        else:
            return slide_text
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": content}],
        temperature=0,
        max_tokens=3000,
    )
    return response.choices[0].message.content or ""


def _parse_json_object(text: str) -> dict[str, Any]:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-zA-Z]*\n?", "", text)
        text = re.sub(r"\n?```$", "", text).strip()
    return json.loads(text)


def _extract_modules(client, model: str, product: str, source: str) -> dict[str, Any]:
    catalog = "\n".join(f"- {mid}: {name}" for mid, name in MODULE_CATALOG.items())
    # 原文过长时截断（qwen3.7-plus 上下文足够大，此处兜底 8 万字）
    prompt = EXTRACT_PROMPT.format(catalog=catalog, product=product, source=source[:80000])
    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        temperature=0,
        max_tokens=8000,
        response_format={"type": "json_object"},
    )
    data = _parse_json_object(response.choices[0].message.content or "{}")
    modules = data.get("modules", data) or {}
    cleaned = {}
    for mid, mod in modules.items():
        mid = str(mid).zfill(2)
        if mid not in MODULE_CATALOG or not isinstance(mod, dict):
            continue
        text = str(mod.get("text") or "").strip()
        if not text:
            continue
        cleaned[mid] = {
            "name": MODULE_CATALOG[mid],
            "text": text,
            "quotes": [str(q) for q in (mod.get("quotes") or []) if str(q).strip()][:8],
        }
    return cleaned


_WS_RE = re.compile(r"[\s，,。.、；;：:！!？?\-—*#|>`]")


def _normalize(text: str) -> str:
    return _WS_RE.sub("", text)


def _verify_quotes(modules: dict[str, Any], source: str) -> None:
    """确定性核验：quote 归一化后是否真是原文子串。不花 LLM 调用。"""
    normalized_source = _normalize(source)
    for mod in modules.values():
        quotes = mod.get("quotes") or []
        hits = sum(1 for q in quotes if _normalize(q) and _normalize(q) in normalized_source)
        mod["quotes_verified"] = hits
        mod["quotes_total"] = len(quotes)
        # verified: 引用全部可在原文找到；unverified 的模块审核时要重点人查
        mod["verified"] = bool(quotes) and hits == len(quotes)


def _run_import(import_id: str) -> None:
    record = get_import(import_id)
    if record is None or record["status"] not in ("queued", "running"):
        return
    workdir = KB_IMPORTS_DIR / import_id
    source_path = next(workdir.glob("source.*"), None)
    if source_path is None:
        raise RuntimeError("源文件丢失")
    _update(import_id, status="running", stage="render", error=None)

    pages = _render_pages(source_path, record["file_ext"])
    _update(import_id, page_count=len(pages), stage="recognize")

    client, route = _vision_client()
    model = _vision_model()
    done_lock = threading.Lock()
    done = {"n": 0}

    def _do(page: dict[str, Any]) -> str:
        text = _recognize_page(client, model, page)
        with done_lock:
            done["n"] += 1
            _update(import_id, pages_done=done["n"])
        return text

    needs_llm = [p for p in pages if p["kind"] != "text" and not (p["kind"] == "slide" and not p.get("images"))]
    with ThreadPoolExecutor(max_workers=RECOGNIZE_CONCURRENCY) as pool:
        page_texts = list(pool.map(_do, pages))
    source_text = "\n\n---\n\n".join(t for t in page_texts if t and t.strip())
    if len(source_text.strip()) < 20:
        raise RuntimeError("识别结果几乎为空，请检查文件内容是否清晰")
    (workdir / "source_text.md").write_text(source_text, encoding="utf-8")

    _update(import_id, stage="extract", source_chars=len(source_text))
    modules = _extract_modules(client, model, record["product_key"], source_text)
    if not modules:
        raise RuntimeError("未能从文档中提取到任何知识模块")

    _update(import_id, stage="verify")
    _verify_quotes(modules, source_text)

    _update(
        import_id,
        status="success",
        stage="done",
        draft_json=json.dumps(
            {"modules": modules, "model": model, "route": route}, ensure_ascii=False
        ),
        finished_at=_now(),
    )


# ---------------------------------------------------------------------------
# apply：审核通过的模块合并进知识库
# ---------------------------------------------------------------------------


def apply_import(
    import_id: str,
    username: str,
    role: str,
    modules: dict[str, str],
    scope: str | None = None,
) -> dict[str, Any]:
    """把用户勾选采纳的模块（module_id → 最终文本，可在审核时手改）合并进知识库。"""
    from . import user_config_store, yaml_store  # noqa: F401 (yaml_store 保留快照先例)

    record = get_import(import_id)
    if record is None:
        raise ValueError("导入任务不存在")
    if record["status"] != "success":
        raise ValueError("任务尚未完成，不能合并")
    if not modules:
        raise ValueError("未选择任何模块")
    scope = scope or record["scope"]
    if scope == "global" and role != "admin":
        raise ValueError("只有管理员能写入全局知识库")

    product_key = record["product_key"]
    if scope == "global":
        kb = user_config_store.load_global_kb()
    else:
        kb = user_config_store.load_effective_kb(username)["data"]

    entry = kb.setdefault(product_key, {"product_name": product_key, "modules": {}})
    entry.setdefault("modules", {})
    applied = []
    for mid, text in modules.items():
        mid = str(mid).zfill(2)
        if mid not in MODULE_CATALOG or not str(text).strip():
            continue
        entry["modules"][mid] = {"name": MODULE_CATALOG[mid], "text": str(text).strip()}
        applied.append(mid)
    if not applied:
        raise ValueError("没有有效的模块内容")

    if scope == "global":
        # 全局是 07 校验的默认标准，写入前快照
        from ..core.paths import RUNS_DIR

        snapshot = RUNS_DIR / f"kb_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        snapshot.write_text(
            json.dumps(user_config_store.load_global_kb(), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        user_config_store.save_global_kb(kb)
    else:
        user_config_store.save_user_kb(username, kb)

    _update(import_id, applied_at=_now())
    return {"applied_modules": applied, "scope": scope, "product_key": product_key}
